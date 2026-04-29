#!/usr/bin/env python3
"""
SARS Project Presentation Generator
Generates a professional PowerPoint presentation for the B.Tech viva
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Color scheme
NAVY = RGBColor(15, 23, 42)  # #0F172A
BLUE = RGBColor(59, 130, 246)  # #3B82F6
CYAN = RGBColor(6, 182, 212)  # #06B6D4
GREEN = RGBColor(34, 197, 94)  # #22C55E
AMBER = RGBColor(245, 158, 11)  # #F59E0B
RED = RGBColor(239, 68, 68)  # #EF4444
WHITE = RGBColor(241, 245, 249)  # #F1F5F9
GRAY = RGBColor(148, 163, 184)  # #94A3B8

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)  # 33.87 cm
prs.slide_height = Inches(7.5)    # 19.05 cm

def add_footer(slide, slide_num, total_slides=19):
    """Add consistent footer to slide"""
    footer_left = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(3), Inches(0.4))
    tf_left = footer_left.text_frame
    tf_left.text = "SARS | SNIST 2024-25"
    tf_left.paragraphs[0].font.size = Pt(10)
    tf_left.paragraphs[0].font.color.rgb = GRAY

    footer_right = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.4))
    tf_right = footer_right.text_frame
    tf_right.text = f"Slide {slide_num} / {total_slides}"
    tf_right.paragraphs[0].font.size = Pt(10)
    tf_right.paragraphs[0].font.color.rgb = GRAY
    tf_right.paragraphs[0].alignment = PP_ALIGN.RIGHT

def set_background(slide, color=NAVY):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_section_label(slide, text):
    """Add section label like '01 — Abstract'"""
    label = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(0.5))
    tf = label.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE

def add_centered_title(slide, text):
    """Add main slide title"""
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.33), Inches(1))
    tf = title.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 1 - TITLE SLIDE
# ========================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_background(slide1)

# Add subtle header line
line_top = slide1.shapes.add_shape(1, Inches(0), Inches(0.2), Inches(13.333), Inches(0.01))
line_top.fill.solid()
line_top.fill.fore_color.rgb = BLUE
line_top.line.color.rgb = BLUE

# Small label
label_sm = slide1.shapes.add_textbox(Inches(4.5), Inches(0.5), Inches(4), Inches(0.3))
tf = label_sm.text_frame
tf.text = "B.Tech Final Year Project — 2024-25"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = CYAN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Main title
title_main = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
tf = title_main.text_frame
tf.word_wrap = True
tf.text = "An Explainable Agentic RAG-Based\nStudent Academic Performance Risk Assessment\nand Advisory System"
for p in tf.paragraphs:
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Acronym
acronym = slide1.shapes.add_textbox(Inches(4), Inches(3.7), Inches(5.333), Inches(0.8))
tf = acronym.text_frame
tf.text = "SARS"
tf.paragraphs[0].font.size = Pt(60)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Student info (left)
student_info = slide1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(5), Inches(1.5))
tf = student_info.text_frame
tf.word_wrap = True
tf.text = "Presented by:\nKethavath Santhosh\nRoll No: 22311A12P9\nB.Tech IV Year — Information Technology"
for i, p in enumerate(tf.paragraphs):
    p.font.size = Pt(12 if i == 0 else 11)
    p.font.color.rgb = WHITE if i == 0 else GRAY
    if i == 0:
        p.font.bold = True

# Guide info (right)
guide_info = slide1.shapes.add_textbox(Inches(7.5), Inches(4.8), Inches(5), Inches(1.5))
tf = guide_info.text_frame
tf.word_wrap = True
tf.text = "Under the Guidance of:\nDr. Sunil Bhutada\nDepartment of Computer Science\nSNIST, Hyderabad"
for i, p in enumerate(tf.paragraphs):
    p.font.size = Pt(12 if i == 0 else 11)
    p.font.color.rgb = WHITE if i == 0 else GRAY
    if i == 0:
        p.font.bold = True

# Bottom line
line_bottom = slide1.shapes.add_shape(1, Inches(0), Inches(7.4), Inches(13.333), Inches(0.01))
line_bottom.fill.solid()
line_bottom.fill.fore_color.rgb = BLUE
line_bottom.line.color.rgb = BLUE

add_footer(slide1, 1)

# ========================
# SLIDE 2 - TABLE OF CONTENTS
# ========================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide2)
add_centered_title(slide2, "Agenda")
add_footer(slide2, 2)

# Content items organized in 2 columns
items_left = [
    ("01", "Abstract"),
    ("02", "Introduction"),
    ("03", "Existing System"),
    ("04", "Proposed System"),
    ("05", "System Requirements"),
    ("06", "System Architecture"),
    ("07", "RAG Pipeline"),
    ("08", "UML Diagrams"),
    ("09", "System Modules"),
]

items_right = [
    ("10", "Implementation"),
    ("11", "Database Design"),
    ("12", "Key Code Snippets"),
    ("13", "Results and Screenshots"),
    ("14", "Performance Evaluation"),
    ("15", "Conclusion"),
    ("16", "Future Scope"),
    ("17", "References"),
]

# Left column
left_col = slide2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5))
tf = left_col.text_frame
tf.word_wrap = True
for num, text in items_left:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"{num}  {text}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.level = 0

# Right column
right_col = slide2.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5.5), Inches(5.5))
tf = right_col.text_frame
tf.word_wrap = True
for num, text in items_right:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"{num}  {text}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.level = 0

# ========================
# SLIDE 3 - ABSTRACT
# ========================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide3)
add_title_section_label(slide3, "01 — Abstract")
add_footer(slide3, 3)

# Main content box
content_box = slide3.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(11.73), Inches(4.2))
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(30, 41, 59)  # Slightly lighter navy
content_box.line.color.rgb = BLUE
content_box.line.width = Pt(2)

# Title in box
title_box = slide3.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(11), Inches(0.5))
tf = title_box.text_frame
tf.text = "What is SARS?"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE

# Body text
body_box = slide3.shapes.add_textbox(Inches(1.2), Inches(1.95), Inches(11), Inches(3))
tf = body_box.text_frame
tf.word_wrap = True
tf.text = ("Academic institutions collect extensive student performance data but lack tools to convert it into early, actionable risk intelligence. "
           "SARS is an end-to-end AI platform that automatically extracts grades from official grade sheet photos using Google Gemini Vision AI, "
           "computes a transparent SARS Risk Score (0–100) based on GPA, backlogs, and attendance, and provides a RAG-based advisory chatbot that "
           "grounds every answer in the student's own academic data with source citations. The system serves both students and faculty through "
           "dedicated dashboards, enabling early intervention before academic risk becomes irreversible.")
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.LEFT

# Bottom highlight boxes
y_start = 5.6
box_width = 3.3
boxes = [
    ("🤖", "AI Vision\nExtraction"),
    ("📊", "Risk Score\n0-100"),
    ("💬", "RAG Advisory\nChat"),
]
x_positions = [1.2, 5.0, 8.8]

for i, (icon, label) in enumerate(boxes):
    box = slide3.shapes.add_shape(1, Inches(x_positions[i]), Inches(y_start), Inches(box_width), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.color.rgb = BLUE
    box.line.width = Pt(2)

    text_box = slide3.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(y_start + 0.15), Inches(box_width - 0.2), Inches(0.7))
    tf = text_box.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 4 - INTRODUCTION
# ========================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide4)
add_title_section_label(slide4, "02 — Introduction")
add_footer(slide4, 4)

# Left half - The Problem
left_box = slide4.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
left_box.line.color.rgb = RED
left_box.line.width = Pt(3)

# Left title
left_title = slide4.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.4))
tf = left_title.text_frame
tf.text = "The Challenge"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = AMBER

# Left bullets
left_content = slide4.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.5))
tf = left_content.text_frame
tf.word_wrap = True
challenges = [
    "Students fail placements due to undetected low CGPA",
    "Faculty manually review grade sheets — slow and error-prone",
    "Risk identified AFTER exams — too late to intervene",
    "No unified view of grades, attendance, and backlogs",
    "Predictive systems give opaque scores — no explanation"
]
for i, ch in enumerate(challenges):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ch
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_before = Pt(4)
    p.level = 0
    p.text = "⚠️  " + p.text

# Right half - Domain Context
right_box = slide4.shapes.add_shape(1, Inches(6.833), Inches(1.2), Inches(6), Inches(5.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
right_box.line.color.rgb = CYAN
right_box.line.width = Pt(3)

# Right title
right_title = slide4.shapes.add_textbox(Inches(7.03), Inches(1.35), Inches(5.6), Inches(0.4))
tf = right_title.text_frame
tf.text = "Domain Context"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = CYAN

# Right bullets
right_content = slide4.shapes.add_textbox(Inches(7.03), Inches(1.85), Inches(5.4), Inches(4.5))
tf = right_content.text_frame
tf.word_wrap = True
domain_items = [
    "Educational Technology & Learning Analytics",
    "Explainable AI, Student Success Analytics",
    "SNIST — JNTUH Affiliated Institution",
    "B.Tech students across all 8 semesters",
    "Placement threshold: 7.5 CGPA minimum at SNIST"
]
for i, item in enumerate(domain_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_before = Pt(4)
    p.level = 0
    p.text = "✅  " + p.text

# Bottom highlight
highlight = slide4.shapes.add_shape(1, Inches(0.5), Inches(7), Inches(12.333), Inches(0.35))
highlight.fill.solid()
highlight.fill.fore_color.rgb = BLUE
highlight.line.width = Pt(0)

highlight_text = slide4.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12), Inches(0.25))
tf = highlight_text.text_frame
tf.text = "Early identification of at-risk students is the single most impactful intervention an institution can make"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 5 - EXISTING SYSTEM
# ========================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide5)
add_title_section_label(slide5, "03 — Existing System")
add_footer(slide5, 5)

# Process flow
flow_items = [
    "Grade Sheets\nPrinted",
    "Faculty Manually\nReviews",
    "Excel/Register\nUpdated",
    "Risk Identified\nLate",
    "Intervention\nToo Late"
]
x_start = 0.8
box_width = 2.2
spacing = 0.2

for i, item in enumerate(flow_items):
    x = x_start + i * (box_width + spacing)
    # Flow box
    flow_box = slide5.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(box_width), Inches(1))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(100, 30, 30)  # Dark red
    flow_box.line.color.rgb = RED
    flow_box.line.width = Pt(2)

    # Text
    text_box = slide5.shapes.add_textbox(Inches(x + 0.1), Inches(1.4), Inches(box_width - 0.2), Inches(0.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.text = item
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow
    if i < len(flow_items) - 1:
        arrow_x = x + box_width + spacing * 0.5
        arrow = slide5.shapes.add_connector(1, Inches(arrow_x), Inches(1.8), Inches(arrow_x + 0.3), Inches(1.8))
        arrow.line.color.rgb = AMBER
        arrow.line.width = Pt(2)

# Limitations grid (2 rows x 3 columns)
limitations = [
    "❌ Manual Process\n— Error prone,\ntime consuming",
    "❌ Late Detection\n— After exams,\nnot before",
    "❌ No Unified View\n— Grades, attendance\nscattered",
    "❌ Black Box\nPredictions\n— No explanations",
    "❌ No Advisory\nSupport\n— Students get\nno guidance",
    "❌ Faculty\nOverload\n— Manual monitoring\nat scale",
]

y_start = 2.6
for i, limit in enumerate(limitations):
    row = i // 3
    col = i % 3
    x = 0.8 + col * (4)
    y = y_start + row * 1.8

    # Limitation card
    card = slide5.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = RED
    card.line.width = Pt(2)

    # Text
    text_box = slide5.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(3.5), Inches(1.3))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.text = limit
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ========================
# SLIDE 6 - PROPOSED SYSTEM
# ========================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide6)
add_title_section_label(slide6, "04 — Proposed System")
add_footer(slide6, 6)

# Solution statement
solution_box = slide6.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(11.73), Inches(0.6))
solution_box.fill.solid()
solution_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
solution_box.line.color.rgb = GREEN
solution_box.line.width = Pt(2)

solution_text = slide6.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11.33), Inches(0.4))
tf = solution_text.text_frame
tf.text = "SARS replaces manual academic monitoring with an intelligent, automated, explainable AI platform that serves both students and faculty in real time."
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Contribution cards
contributions = [
    (BLUE, "🤖 Contribution 1", "AI Vision Grade\nExtraction", "Upload photo → Gemini reads it\n→ Structured data in 10 seconds"),
    (CYAN, "📊 Contribution 2", "Explainable Risk\nScoring", "Credits-weighted CGPA\n→ SARS Score 0-100 with\nfactor breakdown"),
    (GREEN, "💬 Contribution 3", "RAG Advisory\nChat", "Student asks question\n→ Grounded answer +\nsource citations"),
]

y_card = 2.1
for i, (color, title, subtitle, desc) in enumerate(contributions):
    x = 0.8 + i * 4.1

    card = slide6.shapes.add_shape(1, Inches(x), Inches(y_card), Inches(3.9), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Title
    title_box = slide6.shapes.add_textbox(Inches(x + 0.15), Inches(y_card + 0.15), Inches(3.6), Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide6.shapes.add_textbox(Inches(x + 0.15), Inches(y_card + 0.52), Inches(3.6), Inches(0.4))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide6.shapes.add_textbox(Inches(x + 0.15), Inches(y_card + 1), Inches(3.6), Inches(1.1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Advantages vs Existing table
table_y = 4.65
table_box = slide6.shapes.add_shape(1, Inches(0.8), Inches(table_y), Inches(11.73), Inches(2))
table_box.fill.solid()
table_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
table_box.line.color.rgb = BLUE
table_box.line.width = Pt(1)

headers = ["Existing System", "SARS System"]
rows = [
    ("Manual extraction", "AI Vision — automatic"),
    ("Late risk detection", "Real-time after upload"),
    ("Opaque scores", "Full factor explanation"),
    ("No advisory support", "RAG-based personalized chat"),
]

col_width = 5.8
for col, header in enumerate(headers):
    header_box = slide6.shapes.add_textbox(Inches(0.9 + col * col_width), Inches(table_y + 0.1), Inches(col_width - 0.2), Inches(0.35))
    tf = header_box.text_frame
    tf.text = header
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

for row_idx, (existing, sars) in enumerate(rows):
    y_offset = table_y + 0.5 + row_idx * 0.35

    existing_box = slide6.shapes.add_textbox(Inches(0.9), Inches(y_offset), Inches(col_width - 0.2), Inches(0.3))
    tf = existing_box.text_frame
    tf.text = existing
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE

    sars_box = slide6.shapes.add_textbox(Inches(0.9 + col_width), Inches(y_offset), Inches(col_width - 0.2), Inches(0.3))
    tf = sars_box.text_frame
    tf.text = sars
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE

# ========================
# SLIDE 7 - SYSTEM REQUIREMENTS
# ========================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide7)
add_title_section_label(slide7, "05 — System Requirements")
add_footer(slide7, 7)

# Left column - Hardware
hw_box = slide7.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
hw_box.fill.solid()
hw_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
hw_box.line.color.rgb = BLUE
hw_box.line.width = Pt(2)

hw_title = slide7.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.4))
tf = hw_title.text_frame
tf.text = "💻 Hardware Requirements"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE

hw_specs = [
    ("Processor", "Intel Core i5 or higher"),
    ("RAM", "Minimum 8 GB (16 GB recommended)"),
    ("Storage", "20 GB free disk space"),
    ("Network", "Internet connection (Gemini API calls)"),
    ("Display", "1366 x 768 resolution minimum"),
    ("OS", "Windows 10/11 or Ubuntu 20.04+"),
]

hw_content = slide7.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.5))
tf = hw_content.text_frame
tf.word_wrap = True
for i, (label, value) in enumerate(hw_specs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"{label}: {value}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.space_before = Pt(3)

# Right column - Software
sw_box = slide7.shapes.add_shape(1, Inches(6.833), Inches(1.2), Inches(6), Inches(5.5))
sw_box.fill.solid()
sw_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
sw_box.line.color.rgb = CYAN
sw_box.line.width = Pt(2)

sw_title = slide7.shapes.add_textbox(Inches(7.03), Inches(1.35), Inches(5.6), Inches(0.4))
tf = sw_title.text_frame
tf.text = "⚙️  Software Requirements"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = CYAN

# Backend Stack
backend_label = slide7.shapes.add_textbox(Inches(7.03), Inches(1.85), Inches(5.4), Inches(0.25))
tf = backend_label.text_frame
tf.text = "Backend Stack:"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE

backend_items = [
    "Python 3.11, FastAPI 0.111.0",
    "SQLAlchemy 2.0, SQLite",
    "PyMuPDF 1.24.1, Pillow 10.3.0",
    "ChromaDB 0.5.3, python-jose, bcrypt",
]

backend_content = slide7.shapes.add_textbox(Inches(7.13), Inches(2.12), Inches(5.2), Inches(1.3))
tf = backend_content.text_frame
tf.word_wrap = True
for i, item in enumerate(backend_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.space_before = Pt(2)

# Frontend Stack
frontend_label = slide7.shapes.add_textbox(Inches(7.03), Inches(3.55), Inches(5.4), Inches(0.25))
tf = frontend_label.text_frame
tf.text = "Frontend Stack:"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = CYAN

frontend_items = [
    "React 18.2.0, React Router DOM 6.23.0",
    "Axios 1.7.0, react-markdown",
]

frontend_content = slide7.shapes.add_textbox(Inches(7.13), Inches(3.82), Inches(5.2), Inches(0.8))
tf = frontend_content.text_frame
tf.word_wrap = True
for i, item in enumerate(frontend_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.space_before = Pt(2)

# AI Services
ai_label = slide7.shapes.add_textbox(Inches(7.03), Inches(4.75), Inches(5.4), Inches(0.25))
tf = ai_label.text_frame
tf.text = "AI Services:"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN

ai_items = [
    "Google Gemini 2.5 Flash (Vision + Text)",
    "Text-embedding-004 (Embeddings)",
    "API Key: Google AI Studio (free tier)",
]

ai_content = slide7.shapes.add_textbox(Inches(7.13), Inches(5.02), Inches(5.2), Inches(1))
tf = ai_content.text_frame
tf.word_wrap = True
for i, item in enumerate(ai_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.space_before = Pt(2)

# ========================
# SLIDE 8 - SYSTEM ARCHITECTURE (Placeholder)
# ========================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide8)
add_title_section_label(slide8, "06 — System Architecture")
add_footer(slide8, 8)

# Placeholder box
placeholder = slide8.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)
placeholder.line.color.rgb = BLUE
placeholder.line.width = Pt(2)
placeholder.line.dash_style = 3  # Dashed

placeholder_text = slide8.shapes.add_textbox(Inches(1.7), Inches(3.5), Inches(10), Inches(1))
tf = placeholder_text.text_frame
tf.text = "System Architecture Diagram\n(Insert diagram image here)"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].font.color.rgb = BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

caption = slide8.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
tf = caption.text_frame
tf.text = "Five-layer architecture: Users → React Frontend → FastAPI Backend → SQLite Database → Gemini AI Services"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 9 - RAG PIPELINE (Placeholder)
# ========================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide9)
add_title_section_label(slide9, "07 — RAG Pipeline")
add_footer(slide9, 9)

# Placeholder box
placeholder = slide9.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)
placeholder.line.color.rgb = CYAN
placeholder.line.width = Pt(2)
placeholder.line.dash_style = 3  # Dashed

placeholder_text = slide9.shapes.add_textbox(Inches(1.7), Inches(3.5), Inches(10), Inches(1))
tf = placeholder_text.text_frame
tf.text = "RAG Pipeline Diagram\n(Insert diagram image here)"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].font.color.rgb = CYAN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

caption = slide9.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
tf = caption.text_frame
tf.text = "Student data chunked into 7 types → Gemini embeddings → ChromaDB → Semantic retrieval → Grounded response + Source citations"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 10 - UML DIAGRAMS
# ========================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide10)
add_title_section_label(slide10, "08 — UML Diagrams")
add_footer(slide10, 10)

# Simplified UML content
uml_boxes = [
    (0.5, 1.2, "Use Case Diagram", "Students and Teachers\nlogin, upload marks,\nview risks, interactions"),
    (6.766, 1.2, "Activity Diagram", "Upload → Extract\n→ Review → Save\n→ Compute Risk"),
    (0.5, 4, "Sequence Diagram", "Frontend → Backend\n→ Database → Risk Engine\n→ Response Chain"),
    (6.766, 4, "ER Diagram Overview", "Users, Students,\nSemester Records,\nRisk Scores, Chats"),
]

for x, y, title, content in uml_boxes:
    box = slide10.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.color.rgb = BLUE
    box.line.width = Pt(2)

    title_box = slide10.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(5.4), Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN

    content_box = slide10.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(5.4), Inches(1.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = content
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ========================
# SLIDE 11 - SYSTEM MODULES
# ========================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide11)
add_title_section_label(slide11, "09 — System Modules")
add_footer(slide11, 11)

modules = [
    (0.5, 1.2, BLUE, "🔐 Authentication", ["JWT token login", "bcrypt hashing", "Role separation", "Rate limiting"]),
    (4.5, 1.2, CYAN, "📄 Grade Extraction", ["PDF/Image upload", "Image enhancement", "Gemini Vision OCR", "JSON parsing", "User review"]),
    (8.5, 1.2, AMBER, "⚡ Risk Engine", ["CGPA formula", "GPA + Backlog", "Attendance risk", "Risk levels", "Confidence metric"]),
    (0.5, 3.9, BLUE, "🧠 RAG Advisory", ["7 chunk types", "ChromaDB embeddings", "Semantic retrieval", "Grounded LLM", "Source citations"]),
    (4.5, 3.9, GREEN, "👨‍🎓 Student Dashboard", ["Overview cards", "Upload marks", "Risk page", "Attendance", "AI advisory"]),
    (8.5, 3.9, RED, "👩‍🏫 Teacher Dashboard", ["Risk monitor", "Student profile", "Interventions", "Analytics", "CSV export"]),
]

for x, y, color, title, points in modules:
    box = slide11.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.5), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    box.line.color.rgb = color
    box.line.width = Pt(2)

    title_box = slide11.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(3.2), Inches(0.35))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    points_box = slide11.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(3.2), Inches(1.75))
    tf = points_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(8.5)
        p.font.color.rgb = GRAY
        p.space_before = Pt(1)

# ========================
# SLIDE 12 - IMPLEMENTATION
# ========================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide12)
add_title_section_label(slide12, "10 — Implementation")
add_footer(slide12, 12)

# Left side - Implementation Workflow
workflow_steps = [
    "📝 User Registration",
    "📤 Grade Upload",
    "✅ Review Grades",
    "⚡ Risk Computation",
    "🧠 RAG Indexing",
    "💬 AI Advisory",
    "👩‍🏫 Teacher Monitor",
]

workflow_box = slide12.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
tf = workflow_box.text_frame
tf.word_wrap = True
for step in workflow_steps:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_before = Pt(8)

# Right side - Project Structure
structure_text = """sars/
├── backend/
│   ├── app/
│   │   ├── routes/
│   │   │   ├── auth.py
│   │   │   ├── student.py
│   │   │   └── teacher.py
│   │   ├── services/
│   │   │   ├── grade_extractor.py
│   │   │   ├── risk_engine.py
│   │   │   ├── rag_service.py
│   │   │   └── advisor.py
│   │   └── models/
│   │       └── models.py
│   └── config.py
└── frontend/
    └── src/
        ├── pages/student/
        └── pages/teacher/"""

structure_box = slide12.shapes.add_shape(1, Inches(6.5), Inches(1.2), Inches(6.333), Inches(5.5))
structure_box.fill.solid()
structure_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
structure_box.line.color.rgb = BLUE
structure_box.line.width = Pt(1)

structure_content = slide12.shapes.add_textbox(Inches(6.65), Inches(1.35), Inches(6.1), Inches(5.2))
tf = structure_content.text_frame
tf.word_wrap = False
tf.text = structure_text
tf.paragraphs[0].font.name = "Courier New"
tf.paragraphs[0].font.size = Pt(8)
tf.paragraphs[0].font.color.rgb = CYAN

# ========================
# SLIDE 13 - DATABASE DESIGN
# ========================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide13)
add_title_section_label(slide13, "11 — Database Design")
add_footer(slide13, 13)

# Stats cards
stats = [
    ("11 Tables", "Database Schema"),
    ("SQLite +", "SQLAlchemy ORM"),
    ("FK Indexed", "Performance"),
]

for i, (stat, label) in enumerate(stats):
    x = 0.8 + i * 3.8
    card = slide13.shapes.add_shape(1, Inches(x), Inches(0.9), Inches(3.6), Inches(0.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = BLUE
    card.line.width = Pt(1)

    stat_text = slide13.shapes.add_textbox(Inches(x + 0.1), Inches(0.95), Inches(3.4), Inches(0.3))
    tf = stat_text.text_frame
    tf.text = stat
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    label_text = slide13.shapes.add_textbox(Inches(x + 0.1), Inches(1.25), Inches(3.4), Inches(0.2))
    tf = label_text.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Tables
tables = [
    ("users", ["id", "email", "full_name", "password_hash", "role"]),
    ("students", ["id", "user_id(FK)", "branch", "cgpa"]),
    ("semester_records", ["id", "student_id(FK)", "semester", "gpa", "credits"]),
    ("subject_grades", ["id", "semester_id(FK)", "subject", "grade", "credits"]),
    ("attendance_records", ["id", "student_id(FK)", "attended", "total"]),
    ("risk_scores", ["id", "student_id(FK)", "score", "level", "timestamp"]),
]

y_start = 1.9
col_width = 6.2
for i, (table_name, columns) in enumerate(tables):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.5
    y = y_start + row * 1.8

    table_box = slide13.shapes.add_shape(1, Inches(x), Inches(y), Inches(col_width), Inches(1.6))
    table_box.fill.solid()
    table_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    table_box.line.color.rgb = CYAN
    table_box.line.width = Pt(1)

    # Table name
    name_box = slide13.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(col_width - 0.3), Inches(0.25))
    tf = name_box.text_frame
    tf.text = f"📦 {table_name.upper()}"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN

    # Columns
    cols_box = slide13.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.32), Inches(col_width - 0.3), Inches(1.2))
    tf = cols_box.text_frame
    tf.word_wrap = True
    for j, col_name in enumerate(columns):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = col_name
        p.font.size = Pt(8)
        p.font.color.rgb = WHITE
        p.space_before = Pt(1)

# ========================
# SLIDE 14 - KEY CODE SNIPPETS
# ========================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide14)
add_title_section_label(slide14, "12 — Key Code Snippets")
add_footer(slide14, 14)

code_blocks = [
    ("risk_engine.py — SARS Score Formula", """gpa_risk = max(0.0, (7.5 - cgpa) / 7.5 * 100)
backlog_map = {0:0, 1:40, 2:60, 3:80, 4:90}
backlog_risk = backlog_map.get(backlogs, 100)

sars_score = (gpa_risk * 0.40 +
              backlog_risk * 0.35 +
              attendance_risk * 0.25)"""),

    ("risk_engine.py — Credits-Weighted CGPA", """def compute_cgpa(student_id, db):
    records = db.query(SemesterRecord)\\
                .filter_by(student_id=student_id).all()
    total_pts = sum(r.gpa * r.credits
                   for r in records)
    total_cr = sum(r.credits for r in records)
    return round(total_pts / total_cr, 2)"""),

    ("rag_service.py — Semantic Retrieval", """def retrieve(student_id, query, n=4):
    query_emb = genai.embed_content(
        model="text-embedding-004",
        content=query)
    results = db.query(
        query_embeddings=[query_emb],
        n_results=n,
        where={"student_id": str(student_id)})
    return results"""),
]

y_base = 1.2
for i, (label, code) in enumerate(code_blocks):
    y = y_base + i * 1.95

    # Label
    label_box = slide14.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(0.25))
    tf = label_box.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN

    # Code box
    code_box = slide14.shapes.add_shape(1, Inches(0.5), Inches(y + 0.28), Inches(12.333), Inches(1.6))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    code_box.line.color.rgb = BLUE
    code_box.line.width = Pt(1)

    # Code text
    code_text = slide14.shapes.add_textbox(Inches(0.65), Inches(y + 0.38), Inches(12.1), Inches(1.4))
    tf = code_text.text_frame
    tf.word_wrap = False
    tf.text = code
    tf.paragraphs[0].font.name = "Courier New"
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = CYAN

# ========================
# SLIDE 15 - RESULTS AND SCREENSHOTS
# ========================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide15)
add_title_section_label(slide15, "13 — Results and Screenshots")
add_footer(slide15, 15)

screenshots = [
    (0.5, 1.2, "Student Login Page", "Role-based authentication"),
    (4.8, 1.2, "Grade Upload", "Gemini Vision Extraction"),
    (9.1, 1.2, "Risk Score Page", "Factor breakdown card"),
    (0.5, 3.8, "SGPA Bar Chart", "8-semester trend"),
    (4.8, 3.8, "Teacher Risk Monitor", "Class-wide ranking"),
    (9.1, 3.8, "Advisory Chat", "With source citations"),
]

for x, y, title, desc in screenshots:
    img_box = slide15.shapes.add_shape(1, Inches(x), Inches(y), Inches(4), Inches(2.3))
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    img_box.line.color.rgb = BLUE
    img_box.line.width = Pt(2)
    img_box.line.dash_style = 3  # Dashed

    # Placeholder text
    placeholder_txt = slide15.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.85), Inches(3.7), Inches(0.6))
    tf = placeholder_txt.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description
    desc_txt = slide15.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.55), Inches(3.7), Inches(0.6))
    tf = desc_txt.text_frame
    tf.word_wrap = True
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ========================
# SLIDE 16 - PERFORMANCE EVALUATION
# ========================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide16)
add_title_section_label(slide16, "14 — Performance Evaluation")
add_footer(slide16, 16)

# Metrics
metrics = [
    ("100%", "Grade Extraction\nAccuracy", GREEN),
    ("100%", "RAG Retrieval\nAccuracy", GREEN),
    ("73%", "Token Reduction\nvs Baseline", BLUE),
    ("8.55", "Correct CGPA\nMatch", CYAN),
]

for i, (value, label, color) in enumerate(metrics):
    x = 0.8 + i * 3.1
    metric_box = slide16.shapes.add_shape(1, Inches(x), Inches(1), Inches(2.9), Inches(1.2))
    metric_box.fill.solid()
    metric_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    metric_box.line.color.rgb = color
    metric_box.line.width = Pt(2)

    # Value
    val_box = slide16.shapes.add_textbox(Inches(x + 0.1), Inches(1.05), Inches(2.7), Inches(0.4))
    tf = val_box.text_frame
    tf.text = value
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide16.shapes.add_textbox(Inches(x + 0.1), Inches(1.45), Inches(2.7), Inches(0.7))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    tf.text = label
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Results table
results_title = slide16.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.3))
tf = results_title.text_frame
tf.text = "Grade Extraction Validation"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE

# Table headers
headers = ["Student", "Semesters", "Subjects", "Extracted", "Accuracy"]
header_y = 2.85
for i, header in enumerate(headers):
    header_box = slide16.shapes.add_textbox(Inches(0.5 + i * 2.4), Inches(header_y), Inches(2.3), Inches(0.25))
    tf = header_box.text_frame
    tf.text = header
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE

# Table rows
rows_data = [
    ["Student A", "5", "47", "47", "100%"],
    ["Student B", "2", "19", "19", "100%"],
    ["Total", "7", "66", "66", "100%"],
]

for row_idx, row_data in enumerate(rows_data):
    row_y = 3.15 + row_idx * 0.3
    for col_idx, cell_text in enumerate(row_data):
        cell_box = slide16.shapes.add_textbox(Inches(0.5 + col_idx * 2.4), Inches(row_y), Inches(2.3), Inches(0.25))
        tf = cell_box.text_frame
        tf.text = cell_text
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = WHITE

# RAG Table
rag_title = slide16.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.333), Inches(0.3))
tf = rag_title.text_frame
tf.text = "RAG Semantic Retrieval Accuracy"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = CYAN

# RAG headers
rag_headers = ["Query Type", "Expected Chunk", "Retrieved", "Match"]
rag_header_y = 4.45
for i, header in enumerate(rag_headers):
    header_box = slide16.shapes.add_textbox(Inches(0.5 + i * 3), Inches(rag_header_y), Inches(2.9), Inches(0.25))
    tf = header_box.text_frame
    tf.text = header
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN

# RAG rows
rag_rows = [
    ["Risk level", "risk_summary", "risk_summary", "✅"],
    ["Grade query", "semester_N", "semester_N", "✅"],
    ["Placement", "placement_eligible", "placement", "✅"],
    ["What-if", "cgpa_trajectory", "cgpa_trajectory", "✅"],
]

for row_idx, row_data in enumerate(rag_rows):
    row_y = 4.75 + row_idx * 0.25
    for col_idx, cell_text in enumerate(row_data):
        cell_box = slide16.shapes.add_textbox(Inches(0.5 + col_idx * 3), Inches(row_y), Inches(2.9), Inches(0.25))
        tf = cell_box.text_frame
        tf.text = cell_text
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = WHITE

# ========================
# SLIDE 17 - CONCLUSION
# ========================
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide17)
add_title_section_label(slide17, "15 — Conclusion")
add_footer(slide17, 17)

# Left half - What We Achieved
left_box = slide17.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
left_box.line.color.rgb = GREEN
left_box.line.width = Pt(3)

left_title = slide17.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.4))
tf = left_title.text_frame
tf.text = "✅ What We Achieved"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN

achievements = [
    "AI Vision extracts grades with 100% accuracy",
    "SARS Score correctly classifies all risk levels",
    "RAG pipeline reduces tokens by 73%",
    "Source citations make responses explainable",
    "Dual dashboard serves students and faculty",
]

left_content = slide17.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.2))
tf = left_content.text_frame
tf.word_wrap = True
for i, ach in enumerate(achievements):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ach
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_before = Pt(5)
    p.level = 0
    p.text = "🎯 " + p.text

# Right half - Future Scope
right_box = slide17.shapes.add_shape(1, Inches(6.833), Inches(1.2), Inches(6), Inches(5.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
right_box.line.color.rgb = BLUE
right_box.line.width = Pt(3)

right_title = slide17.shapes.add_textbox(Inches(7.03), Inches(1.35), Inches(5.6), Inches(0.4))
tf = right_title.text_frame
tf.text = "🚀 Future Enhancements"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLUE

future_items = [
    "PostgreSQL — production-grade database",
    "Mobile App — React Native camera upload",
    "Multi-language — Hindi and regional languages",
    "LMS Integration — direct portal API connectors",
    "Advanced Agentic Workflows — LangGraph integration",
    "Longitudinal Study — validate across institutions",
    "Email Alerts — notify HIGH risk students",
]

right_content = slide17.shapes.add_textbox(Inches(7.03), Inches(1.85), Inches(5.4), Inches(4.5))
tf = right_content.text_frame
tf.word_wrap = True
for i, item in enumerate(future_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.space_before = Pt(3)
    p.level = 0

# ========================
# SLIDE 18 - REFERENCES
# ========================
slide18 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide18)
add_title_section_label(slide18, "16 — References")
add_footer(slide18, 18)

references = [
    '[1] R. S. J. D. Baker and P. S. Inventado, "Educational data mining and learning analytics," Springer, 2014.',
    '[2] A. Essa and H. Ayad, "Improving student success using predictive models," Research in Learning Technology, vol. 20, 2012.',
    '[3] K. E. Arnold and M. D. Pistilli, "Course signals at Purdue," Proc. LAK 2012, pp. 267-270.',
    '[4] S. M. Jayaprakash et al., "Early alert of academically at-risk students," Journal of Learning Analytics, vol. 1, no. 1, 2014.',
    '[5] P. Lewis et al., "Retrieval-augmented generation for knowledge-intensive NLP tasks," NeurIPS vol. 33, 2020.',
    '[6] Y. Gao et al., "Retrieval-augmented generation for large language models: A survey," arXiv:2312.10997, 2023.',
    '[7] H. Khosravi et al., "Explainable artificial intelligence in education," Computers and Education: AI, vol. 3, 2022.',
    '[8] E. Kasneci et al., "ChatGPT for good? On opportunities and challenges of LLMs for education," Learning and Individual Differences, vol. 103, 2023.',
    '[9] Google AI, "Gemini API Documentation," https://ai.google.dev, 2024.',
    '[10] FastAPI Documentation, https://fastapi.tiangolo.com, 2024.',
]

refs_box = slide18.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
tf = refs_box.text_frame
tf.word_wrap = True
for i, ref in enumerate(references):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ref
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.space_before = Pt(2)
    p.level = 0

add_footer(slide18, 18)

# ========================
# SLIDE 19 - THANK YOU
# ========================
slide19 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide19)

# Add footer lines
line_top = slide19.shapes.add_shape(1, Inches(0), Inches(0.2), Inches(13.333), Inches(0.01))
line_top.fill.solid()
line_top.fill.fore_color.rgb = BLUE
line_top.line.color.rgb = BLUE

line_bottom = slide19.shapes.add_shape(1, Inches(0), Inches(7.4), Inches(13.333), Inches(0.01))
line_bottom.fill.solid()
line_bottom.fill.fore_color.rgb = BLUE
line_bottom.line.color.rgb = BLUE

# Thank You text
thankyou_txt = slide19.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(0.8))
tf = thankyou_txt.text_frame
tf.text = "Thank You"
tf.paragraphs[0].font.size = Pt(62)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Q&A
qa_txt = slide19.shapes.add_textbox(Inches(1), Inches(3.1), Inches(11.333), Inches(0.6))
tf = qa_txt.text_frame
tf.text = "Questions & Answers"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.color.rgb = CYAN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Student info
info_txt = slide19.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1.5))
tf = info_txt.text_frame
tf.word_wrap = True
tf.text = ("Kethavath Santhosh | 22311A12P9\n"
          "Department of Information Technology\n"
          "SNIST, Hyderabad | 2024-25\n"
          "Guide: Dr. Sunil Bhutada")
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
for p in tf.paragraphs:
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

# Bottom line
bottom_txt = slide19.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.6))
tf = bottom_txt.text_frame
tf.text = "SARS — Student Academic Performance Risk Assessment and Advisory System"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_footer(slide19, 19)

# Save presentation
output_path = Path("c:/Users/hp/Downloads/sars_goal1/SARS_Viva_Presentation.pptx")
prs.save(str(output_path))
print(f"[SUCCESS] Presentation saved to: {output_path}")
print(f"[INFO] Total slides: 19")
print(f"[INFO] Color scheme: Dark Navy (#0F172A) with Electric Blue accents")
print(f"[INFO] Format: 16:9 Widescreen (13.333\" x 7.5\")")
